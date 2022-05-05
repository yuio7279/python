from distutils import text_file
from turtle import width
from pptx import Presentation #라이브러리
from pptx.util import Inches #사진, 표등을 그리기 위해

prs = Presentation() #파워포인트 선언

# for i in range(0, 11):
#     title_slide_layout = prs.slide_layouts[i]   #슬라이드 종류 선택
#     slide = prs.slides.add_slide(title_slide_layout) #슬라이드 추가

# prs.save('add all slides.pptx')

# for i in range(0, 11):
#     print("-----------[%d] ---------" %(i))
#     slide = prs.slides.add_slide(prs.slide_layouts[i])
#     for shape in slide.placeholders:
#         print('%d %s' %(shape.placeholder_format.idx, shape.name))


title_slide_layout = prs.slide_layouts[0] #제목슬라이드
slide = prs.slides.add_slide(title_slide_layout)   #슬라이드 추가

# 제목 - 제목에 값넣기
title = slide.placeholders[0] #제목
title.text = "Hello, World!" 

#부제목
subtitle = slide.placeholders[1] # 부재목상자는 [1]번지
subtitle.text = "python-pptx was here!"

#저장
# prs.save('test.pptx')

bullet_slide_layout = prs.slide_layouts[1] #제목 및 내용 슬라이드
slide = prs.slides.add_slide(bullet_slide_layout) #기존에 있던 슬라이드에 추가

#제목
title_shape = slide.placeholders[0]
title_shape.text = 'Adding a Bullet Slide'

#내용
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = 'Find the bullet slide layout'

#단락 추가
p = tf.add_paragraph()
p.text = 'Use _TexxtFrame.text for first bullet'
p.level = 1  #1 : 들여쓰기 레벨

#단락 추가
p = tf.add_paragraph()
p.texxt = 'Use _TextFrame.add_paragraph() for subsequent bullets'
p.level = 2 #2 : 들여쓰기 레벨

# prs.save('text.pptx')

img_path = 'C:\\Users\super\Desktop\\2-B_KYS\\python_workspace\\0504\\res\\powerpoint_hanling\\slide_test.jpg'

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

left = top = Inches(1)
width = height= Inches(1)
#width heigth가 없는경우 원본 사이즈로
pic = slide.shapes.add_picture(img_path,left,top,width=width,height=height)
# ,left,top,width=width,heigth=heigth
left = Inches(3)
width = Inches(5.5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path,left,top,width=width,height=height)

# prs.save('test.pptx')

title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
shapes = slide.shapes

title_shape = slide.placeholders[0]
title_shape.text = "Adding a Table"

rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

#set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)

#write column headings
table.cell(0,0).text=  "Foo"
table.cell(0,1).text = 'Bar'

#write body cells
table.cell(1,0).text = 'Baz'
table.cell(1,1).text = 'Qux'

prs.save('test.pptx')
